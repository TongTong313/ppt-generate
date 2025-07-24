#!/usr/bin/env python3

import http.server
import io
import os
import re
import requests
import urllib.parse
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import cssutils

SHORT_TEXT_LIMIT_CHARS = 75
TITLE_FONT_PT = 44
SUBTITLE_FONT_PT = 28
CONTENT_FONT_PT = 18
SLIDE_BLANK_LAYOUT = 6
# 16:9 比例的幻灯片尺寸
SLIDE_WIDTH_INCHES = 13.33
SLIDE_HEIGHT_INCHES = 7.5
SLIDE_SMALL_MARGIN_INCHES = 0.5
COLUMN_MARGIN_INCHES = 0.2
HEIGHT_MARGIN_INCHES = 0.2

# Init configuration default values
debug_logs = False
debug_slides = False
server_port = 8080

# Parse configuration environment variables
html2pptx_debug_logs = os.getenv("HTML2PPTX_DEBUG_LOGS", "false")
html2pptx_debug_slides = os.getenv("HTML2PPTX_DEBUG_SLIDES", "false")
html2pptx_server_port = os.getenv("HTML2PPTX_PORT", "8080")

# Set configuration variables depending on environment variables
if html2pptx_debug_logs.lower() == "true":
    debug_logs = True
if html2pptx_debug_slides.lower() == "true":
    debug_slides = True
server_port = int(html2pptx_server_port)

# Print configuration
print("Configuration:")
print("debug_logs:", debug_logs)
print("debug_slides:", debug_slides)
print("server_port:", server_port)


def html_to_pptx(html_content, output_file):
    # 直接使用传入的HTML内容字符串
    css_selector = ".slide-container"  # 选择幻灯片容器
    slides = html_to_slides(html_content, css_selector)
    prs_bytes_stream = slides_to_pptx(slides)

    # 保存PPT文件
    with open(output_file, "wb") as f:
        f.write(prs_bytes_stream.getvalue())

    return prs_bytes_stream


def html_to_slides(html_string, css_selector, style_map=None):
    soup = BeautifulSoup(html_string, "html.parser")
    useful_content = soup.select(css_selector)
    slides = []
    for slide_container in useful_content:
        # 每个slide_container整体作为一页slide
        slide_content = html_to_slide(slide_container, style_map)
        slides.append(slide_content)
    return slides


def html_to_slide(parent_tag, style_map=None):
    return parse_tag_contents_with_style(parent_tag, style_map)


# 单位换算常量
REM_PX = 16  # 1rem=16px
PX_PT = 0.75  # 1px=0.75pt
PPT_WIDTH_PX = 1280  # 13.33英寸
PPT_HEIGHT_PX = 720  # 7.5英寸

import math


def css_value_to_pt(val,
                    base_font_px=REM_PX,
                    ppt_width_px=PPT_WIDTH_PX,
                    ppt_height_px=PPT_HEIGHT_PX):
    """
    将CSS长度单位（rem, em, px, pt, vw, vh, %）统一转换为pt
    """
    if isinstance(val, (int, float)):
        return float(val)
    if not isinstance(val, str):
        return 0
    val = val.strip().lower()
    try:
        if val.endswith('pt'):
            return float(val[:-2])
        elif val.endswith('px'):
            return float(val[:-2]) * PX_PT
        elif val.endswith('rem'):
            return float(val[:-3]) * REM_PX * PX_PT
        elif val.endswith('em'):
            return float(val[:-2]) * base_font_px * PX_PT
        elif val.endswith('vw'):
            return float(val[:-2]) * ppt_width_px / 100 * PX_PT
        elif val.endswith('vh'):
            return float(val[:-2]) * ppt_height_px / 100 * PX_PT
        elif val.endswith('%'):
            return float(val[:-1]) / 100 * base_font_px * PX_PT
        else:
            return float(val)
    except Exception:
        return 0


# 修改parse_css_styles和get_element_style，所有涉及font-size、margin、padding、line-height等属性都用css_value_to_pt转换
# 在fill_slide中，所有涉及pt的地方都用css_value_to_pt


def parse_css_styles(css_content):
    """
    解析CSS内容，返回样式映射表：{selector: {property: value, ...}}
    """
    sheet = cssutils.parseString(css_content)
    style_map = {}
    for rule in sheet:
        if rule.type == rule.STYLE_RULE:
            selector = rule.selectorText
            props = {}
            for item in rule.style:
                props[item.name] = item.value
            style_map[selector] = props
    return style_map


def get_element_style(tag, style_map):
    """
    根据tag的class、id、标签名，合并所有相关样式，返回最终样式dict
    优先级：id > class > tag
    """
    styles = {}
    # 标签名
    tag_selector = tag.name
    if tag_selector and f"{tag_selector}" in style_map:
        styles.update(style_map[f"{tag_selector}"])
    # class
    if tag.has_attr('class'):
        for cls in tag['class']:
            class_selector = f".{cls}"
            if class_selector in style_map:
                styles.update(style_map[class_selector])
    # id
    if tag.has_attr('id'):
        id_selector = f"#{tag['id']}"
        if id_selector in style_map:
            styles.update(style_map[id_selector])
    # style属性内联
    if tag.has_attr('style'):
        inline = cssutils.parseStyle(tag['style'])
        for item in inline:
            styles[item.name] = item.value
    return styles


def parse_tag_contents_with_style(tag, style_map=None):
    tag_data = []
    for children_content_tag in tag.children:
        if children_content_tag.name is not None:
            style = get_element_style(children_content_tag,
                                      style_map) if style_map else {}
            if children_content_tag.name == "img":
                tag_data.append({
                    "type": "img",
                    "src": children_content_tag["src"],
                    "style": style
                })
            elif children_content_tag.name == "h1":
                text = children_content_tag.get_text().strip()
                if text:
                    tag_data.append({
                        "type": "title",
                        "text": text,
                        "style": style
                    })
            elif children_content_tag.name == "h2":
                text = children_content_tag.get_text().strip()
                if text:
                    tag_data.append({
                        "type": "subtitle",
                        "text": text,
                        "style": style
                    })
            elif children_content_tag.name == "ul":
                for li in children_content_tag.find_all("li", recursive=False):
                    li_style = get_element_style(
                        li, style_map) if style_map else {}
                    text = li.get_text().strip()
                    if text:
                        tag_data.append({
                            "type": "bullet",
                            "text": text,
                            "style": li_style
                        })
            elif children_content_tag.name == "ol":
                idx = 1
                for li in children_content_tag.find_all("li", recursive=False):
                    li_style = get_element_style(
                        li, style_map) if style_map else {}
                    text = li.get_text().strip()
                    if text:
                        tag_data.append({
                            "type": "numbered",
                            "text": text,
                            "number": idx,
                            "style": li_style
                        })
                        idx += 1
            elif children_content_tag.name in ["p", "div"]:
                text = children_content_tag.get_text().strip()
                if text:
                    tag_data.append({
                        "type": "text",
                        "text": text,
                        "style": style
                    })
            elif children_content_tag.string is not None:
                if children_content_tag.string.strip() != "":
                    tag_data.append({
                        "type": "text",
                        "text": children_content_tag.string.strip(),
                        "style": style
                    })
            else:
                tag_data.extend(
                    parse_tag_contents_with_style(children_content_tag,
                                                  style_map))
    return tag_data


def slides_to_pptx(slides):
    prs = Presentation()

    # 设置16:9幻灯片尺寸
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    for slide in slides:
        if debug_logs:
            print(
                "============================================ NEW SLIDE ============================================"
            )
        fill_slide(prs, slide)
        if debug_logs:
            print(
                "============================================ END SLIDE ============================================"
            )
    prs_bytes_stream = io.BytesIO()
    prs.save(prs_bytes_stream)

    # Test to save bytes stream directly to file
    # Uncomment if needed
    # with open("test.pptx", 'wb') as out:
    #     out.write(prs_bytes_stream.getvalue())
    return prs_bytes_stream


# 修改fill_slide，支持新结构和样式应用


def fill_slide(prs, slide):
    if not slide:
        return
    prs_slide_layout = prs.slide_layouts[6]
    prs_slide = prs.slides.add_slide(prs_slide_layout)
    background = prs_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(248, 249, 250)
    top_bar = prs_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                         Inches(0), Inches(SLIDE_WIDTH_INCHES),
                                         Inches(0.067))
    top_bar_fill = top_bar.fill
    top_bar_fill.solid()
    top_bar_fill.fore_color.rgb = RGBColor(102, 126, 234)
    top_bar.line.fill.background()
    circle = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(SLIDE_WIDTH_INCHES - 1.2),
                                        Inches(SLIDE_HEIGHT_INCHES - 1.2),
                                        Inches(0.5), Inches(0.5))
    circle_fill = circle.fill
    circle_fill.solid()
    circle_fill.fore_color.rgb = RGBColor(102, 126, 234)
    circle_fill.transparency = 0.9
    circle.line.fill.background()
    margin_top = Inches(0.79)  # 60px
    margin_left = Inches(1.11)  # 80px
    margin_right = Inches(1.11)  # 80px
    margin_bottom = Inches(0.42)  # 30px
    available_width = Inches(SLIDE_WIDTH_INCHES) - margin_left - margin_right
    current_top = margin_top
    # 分类内容
    titles = []
    subtitles = []
    bullets = []
    numbers = []
    texts = []
    images = []
    for item in slide:
        if item["type"] == "title":
            titles.append(item)
        elif item["type"] == "subtitle":
            subtitles.append(item)
        elif item["type"] == "bullet":
            bullets.append(item)
        elif item["type"] == "numbered":
            numbers.append(item)
        elif item["type"] == "text":
            texts.append(item)
        elif item["type"] == "img":
            images.append(item)
    # 渲染标题
    if titles:
        title_text = '\n'.join([t["text"] for t in titles])
        title_style = titles[0]["style"] if titles[0].get("style") else {}
        title_height = Inches(1.8)
        text_box = prs_slide.shapes.add_textbox(margin_left, current_top,
                                                available_width, title_height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = title_text
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.runs[0]
        run.font.size = Pt(css_value_to_pt(title_style.get("font-size", 44)))
        run.font.name = title_style.get("font-family", 'Microsoft YaHei')
        run.font.bold = True
        if "color" in title_style:
            rgb = parse_css_color(title_style["color"])
            if rgb:
                run.font.color.rgb = rgb
        else:
            run.font.color.rgb = RGBColor(44, 62, 80)
        paragraph.line_spacing = css_value_to_pt(
            title_style.get("line-height", 1.2))
        current_top += title_height + Inches(0.32)
    # 渲染副标题
    if subtitles:
        subtitle_text = '\n'.join([t["text"] for t in subtitles])
        subtitle_style = subtitles[0]["style"] if subtitles[0].get(
            "style") else {}
        subtitle_height = Inches(1.0)
        text_box = prs_slide.shapes.add_textbox(margin_left, current_top,
                                                available_width,
                                                subtitle_height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = subtitle_text
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.runs[0]
        run.font.size = Pt(css_value_to_pt(subtitle_style.get("font-size",
                                                              28)))
        run.font.name = subtitle_style.get("font-family", 'Microsoft YaHei')
        if "color" in subtitle_style:
            rgb = parse_css_color(subtitle_style["color"])
            if rgb:
                run.font.color.rgb = rgb
        else:
            run.font.color.rgb = RGBColor(127, 140, 141)
        paragraph.line_spacing = css_value_to_pt(
            subtitle_style.get("line-height", 1.3))
        current_top += subtitle_height + Inches(0.32)
    # 内容区
    content_area_top = current_top
    content_area_height = Inches(
        SLIDE_HEIGHT_INCHES) - current_top - margin_bottom
    # 渲染普通文本
    if texts and not bullets and not numbers:
        text_text = '\n'.join([t["text"] for t in texts])
        text_style = texts[0]["style"] if texts[0].get("style") else {}
        text_box = prs_slide.shapes.add_textbox(margin_left, content_area_top,
                                                available_width,
                                                content_area_height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text_text
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = paragraph.runs[0]
        run.font.size = Pt(css_value_to_pt(text_style.get("font-size", 20)))
        run.font.name = text_style.get("font-family", 'Microsoft YaHei')
        if "color" in text_style:
            rgb = parse_css_color(text_style["color"])
            if rgb:
                run.font.color.rgb = rgb
        else:
            run.font.color.rgb = RGBColor(44, 62, 80)
        paragraph.line_spacing = css_value_to_pt(
            text_style.get("line-height", 1.8))
    # 渲染无序列表
    if bullets:
        list_width = min(available_width, Inches(11.11))
        list_left = margin_left + (available_width - list_width) / 2
        bullet_line_height = 0.7
        list_height = min(content_area_height,
                          Inches(len(bullets) * bullet_line_height + 0.5))
        text_box = prs_slide.shapes.add_textbox(list_left, content_area_top,
                                                list_width, list_height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = Inches(0.42)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        for idx, bullet in enumerate(bullets):
            style = bullet["style"] if bullet.get("style") else {}
            paragraph = text_frame.add_paragraph(
            ) if idx > 0 else text_frame.paragraphs[0]
            bullet_type = style.get("list-style-type", "disc")
            if bullet_type == "disc":
                symbol = "• "
            elif bullet_type == "circle":
                symbol = "◦ "
            elif bullet_type == "triangle":
                symbol = "▶ "
            elif bullet_type == "square":
                symbol = "▪ "
            else:
                symbol = "• "
            symbol_run = paragraph.add_run()
            symbol_run.text = symbol
            symbol_run.font.size = Pt(
                css_value_to_pt(style.get("font-size", 20)))
            symbol_run.font.name = style.get("font-family", 'Microsoft YaHei')
            if "color" in style:
                rgb = parse_css_color(style["color"])
                if rgb:
                    symbol_run.font.color.rgb = rgb
            else:
                symbol_run.font.color.rgb = RGBColor(102, 126, 234)
            text_run = paragraph.add_run()
            text_run.text = bullet["text"]
            text_run.font.size = Pt(css_value_to_pt(style.get("font-size",
                                                              20)))
            text_run.font.name = style.get("font-family", 'Microsoft YaHei')
            if "color" in style:
                rgb = parse_css_color(style["color"])
                if rgb:
                    text_run.font.color.rgb = rgb
            else:
                text_run.font.color.rgb = RGBColor(52, 73, 94)
            paragraph.line_spacing = css_value_to_pt(
                style.get("line-height", 1.6))
            paragraph.space_after = Pt(8)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    # 渲染有序列表
    if numbers:
        list_width = min(available_width, Inches(11.11))
        list_left = margin_left + (available_width - list_width) / 2
        number_line_height = 0.7
        list_height = min(content_area_height,
                          Inches(len(numbers) * number_line_height + 0.5))
        text_box = prs_slide.shapes.add_textbox(list_left, content_area_top,
                                                list_width, list_height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = Inches(0.42)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        for idx, number in enumerate(numbers):
            style = number["style"] if number.get("style") else {}
            paragraph = text_frame.add_paragraph(
            ) if idx > 0 else text_frame.paragraphs[0]
            symbol_run = paragraph.add_run()
            symbol_run.text = f"{number['number']}. "
            symbol_run.font.size = Pt(
                css_value_to_pt(style.get("font-size", 20)))
            symbol_run.font.name = style.get("font-family", 'Microsoft YaHei')
            if "color" in style:
                rgb = parse_css_color(style["color"])
                if rgb:
                    symbol_run.font.color.rgb = rgb
            else:
                symbol_run.font.color.rgb = RGBColor(102, 126, 234)
            text_run = paragraph.add_run()
            text_run.text = number["text"]
            text_run.font.size = Pt(css_value_to_pt(style.get("font-size",
                                                              20)))
            text_run.font.name = style.get("font-family", 'Microsoft YaHei')
            if "color" in style:
                rgb = parse_css_color(style["color"])
                if rgb:
                    text_run.font.color.rgb = rgb
            else:
                text_run.font.color.rgb = RGBColor(52, 73, 94)
            paragraph.line_spacing = css_value_to_pt(
                style.get("line-height", 1.6))
            paragraph.space_after = Pt(8)
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


# 工具函数：CSS颜色转RGBColor


def parse_css_color(color_str):
    if color_str.startswith('#'):
        hex_color = color_str.lstrip('#')
        if len(hex_color) == 3:
            hex_color = ''.join([c * 2 for c in hex_color])
        if len(hex_color) == 6:
            r, g, b = tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
            return RGBColor(r, g, b)
    elif color_str.startswith('rgb'):
        nums = [int(x) for x in re.findall(r'\d+', color_str)]
        if len(nums) >= 3:
            return RGBColor(nums[0], nums[1], nums[2])
    return None


def css_preprocess(css_content):
    """
    过滤掉cssutils不支持的属性和规则，减少报错。
    保留color、font-size、font-family、background、background-color、margin、padding、text-align、line-height、list-style-type等。
    """
    # 移除@keyframes和animation相关
    css_content = re.sub(r'@keyframes[\s\S]+?\}',
                         '',
                         css_content,
                         flags=re.MULTILINE)
    css_content = re.sub(r'animation[^;]*;?', '', css_content)
    # 移除linear-gradient、flex、vw、vh、rem等属性
    css_content = re.sub(r'background\s*:\s*linear-gradient\([^;]+;?', '',
                         css_content)
    css_content = re.sub(r'display\s*:\s*flex;?', '', css_content)
    css_content = re.sub(r'flex(-direction)?\s*:[^;]+;?', '', css_content)
    css_content = re.sub(r'justify-content\s*:[^;]+;?', '', css_content)
    css_content = re.sub(r'align-items\s*:[^;]+;?', '', css_content)
    css_content = re.sub(r'background-clip\s*:[^;]+;?', '', css_content)
    css_content = re.sub(r'-webkit-background-clip\s*:[^;]+;?', '',
                         css_content)
    css_content = re.sub(r'-webkit-text-fill-color\s*:[^;]+;?', '',
                         css_content)
    # 移除vw/vh/rem单位的宽高，保留font-size等常用属性
    css_content = re.sub(
        r'(width|height|max-width|max-height|min-width|min-height)\s*:[^;]*(vw|vh|rem)[^;]*;?',
        '', css_content)
    return css_content


class Html2pptx(http.server.BaseHTTPRequestHandler):

    def do_GET(self):
        # Handle GET requests

        # Set headers
        self.send_response(200)
        self.send_header("Content-type", "text/html")
        self.end_headers()

        # Send index.html page
        with open("index.html", "rb") as out:
            self.wfile.write(out.read())

    def do_POST(self):
        # Handle POST requests

        # Retrieve and decode POST query data
        content_length = int(self.headers["Content-Length"])
        post_data = self.rfile.read(content_length)
        decoded_post_data = urllib.parse.parse_qs(post_data.decode("utf-8"))
        if debug_logs:
            print('decoded_post_data["url"][0]', decoded_post_data["url"][0])
            print('decoded_post_data["selector"][0]',
                  decoded_post_data["selector"][0])

        # Translate HTML to PPTX, retrieves presentation bytes stream
        prs_bytes_stream = html_to_pptx(decoded_post_data["url"][0],
                                        decoded_post_data["selector"][0])

        # Set headers to download the PPTX file
        self.send_response(200)
        self.send_header("Content-Type", "application/octet-stream")
        self.send_header("Content-Disposition",
                         'attachment; filename="presentation.pptx"')
        # This is some unused example code which uses content-length header when transferring a file
        # Since it seems to work here, we won't use it, but the code below will stay
        # here in case we need to use and modify it
        # Source:
        # https://stackoverflow.com/questions/18543640/how-would-i-create-a-python-web-server-that-downloads-a-file-on-any-get-request
        # fs = os.fstat(f.fileno())
        # self.send_header("Content-Length", str(fs.st_size))
        self.end_headers()

        # Send the PPTX presentation
        # Use getvalue() instead of read() with BytesIO to avoid problems
        # Source:
        # https://stackoverflow.com/questions/46981529/why-does-saving-a-presentation-to-a-file-like-object-produce-a-blank-presentatio?noredirect=1&lq=1
        self.wfile.write(prs_bytes_stream.getvalue())


# Setup and start HTTP server with custom Html2pptx handler
# server_address = ("", server_port)
# httpd = http.server.HTTPServer(server_address, Html2pptx)
# print("Serving at port:", server_port)
# httpd.serve_forever()

if __name__ == "__main__":
    # 读取本地HTML文件
    html_file_path = "./test/ppt_test/ppt-demo.html"
    css_file_path = "./test/ppt_test/ppt-template.css"
    output_file = "./test/ppt_test/presentation.pptx"

    try:
        # 1. 读取原始HTML内容
        with open(html_file_path, "r", encoding="utf-8") as file:
            html_content = file.read()

        # 2. 读取PPT友好型CSS
        with open(css_file_path, "r", encoding="utf-8") as file:
            ppt_css_content = file.read()

        # 3. 提取所有.slide-container内容
        soup = BeautifulSoup(html_content, "html.parser")
        slide_elements = soup.select(".slide-container")
        if not slide_elements:
            print("未找到幻灯片容器元素")
            exit(1)

        # 4. 拼接成临时HTML（只包含内容和PPT友好型CSS）
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset=\"UTF-8\">
            <style>
        {ppt_css_content}
            </style>
        </head>
        <body>
        """
        for slide in slide_elements:
            full_html += str(slide) + "\n"
        full_html += "</body></html>"

        # 5. 用该HTML生成PPT
        html_to_pptx(full_html, output_file)
        print(f"PPT已生成：{output_file}")

    except FileNotFoundError as e:
        print(f"文件未找到: {e}")
    except Exception as e:
        print(f"生成PPT时出错：{e}")
