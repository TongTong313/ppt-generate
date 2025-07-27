#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
修复版HTML转PPT转换器
解决生成多页、背景缺失、文本排版问题
"""

import os
import re
from typing import Dict, List, Any, Optional, Tuple
from bs4 import BeautifulSoup, Tag, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import cssutils
from dataclasses import dataclass


@dataclass
class TextStyle:
    """文本样式配置

    Args:
        font_size (int): 字体大小。
        font_color (str): 字体颜色，默认值#333333（深灰色）。
        bold (bool): 是否加粗，默认值False。
        italic (bool): 是否斜体，默认值False。
        alignment (str): 对齐方式，默认值left。
        line_spacing (float): 行间距，默认值1.15。
        space_before (int): 段落前间距，默认值0。
        space_after (int): 段落后间距，默认值4。
    """

    font_size: int
    font_color: str = "#333333"
    bold: bool = False
    italic: bool = False
    alignment: str = "left"
    line_spacing: float = 1.15
    space_before: int = 0
    space_after: int = 4


@dataclass
class SlideConfig:
    """幻灯片配置

    Args:
        width_inches (float): 幻灯片宽度，默认值13.33英寸（16:9 比例）。
        height_inches (float): 幻灯片高度，默认值7.5英寸。
        width_px (int): 幻灯片宽度，默认值1280像素。
        height_px (int): 幻灯片高度，默认值720像素。
        padding_top (float): 顶部内边距，默认值0.8英寸。
        padding_bottom (float): 底部内边距，默认值0.6英寸。
        padding_left (float): 左侧内边距，默认值1.0英寸。
        padding_right (float): 右侧内边距，默认值1.0英寸。
    """

    width_inches: float = 13.33  # 16:9 比例
    height_inches: float = 7.5
    width_px: int = 1280
    height_px: int = 720
    padding_top: float = 0.8
    padding_bottom: float = 0.6
    padding_left: float = 1.0
    padding_right: float = 1.0
    default_font: str = "Microsoft YaHei"

    def __post_init__(self):
        """初始化文本样式配置，这里要针对不同的模板定制化输出，这个字体未必和html是保持一致的

        Attributes:
            text_styles (Dict[str, TextStyle]): 文本样式配置字典，键为样式名称，值为TextStyle实例，具体包括：
            main_title: 主标题
            slide_title: 幻灯片标题
            subtitle: 副标题
            heading: 标题
            body_text: 正文
            list_item: 列表项
            caption: 脚注
            author: 作者
        """
        self.text_styles: Dict[str, TextStyle] = {
            "main_title": TextStyle(
                font_size=36,
                font_color="#2c5aa0",
                bold=True,
                alignment="center",
                space_after=8,
            ),
            "slide_title": TextStyle(
                font_size=28,
                font_color="#2c5aa0",
                bold=True,
                alignment="left",
                space_after=6,
            ),
            "subtitle": TextStyle(
                font_size=20,
                font_color="#666666",
                bold=False,
                alignment="center",
                space_after=6,
            ),
            "heading": TextStyle(
                font_size=18,
                font_color="#333333",
                bold=True,
                alignment="left",
                space_after=4,
            ),
            "body_text": TextStyle(
                font_size=14,
                font_color="#555555",
                bold=False,
                alignment="left",
                space_after=4,
            ),
            "list_item": TextStyle(
                font_size=14,
                font_color="#555555",
                bold=False,
                alignment="left",
                space_after=3,
            ),
            "caption": TextStyle(
                font_size=12,
                font_color="#888888",
                bold=False,
                alignment="center",
                italic=True,
                space_after=3,
            ),
            "author": TextStyle(
                font_size=16,
                font_color="#888888",
                bold=False,
                alignment="center",
                space_after=0,
            ),
        }


class ColorHelper:
    """颜色处理助手"""

    @staticmethod
    def parse_color(color_str: str) -> Optional[RGBColor]:
        """解析颜色字符串
        Args:
            color_str (str): 颜色字符串，支持十六进制、RGB、命名颜色三种方式。
        Returns:
            Optional[RGBColor]: 解析后的RGBColor实例，如果无法解析则返回None。
        """
        if not color_str:
            return None

        color_str = color_str.strip().lower()

        # 1. 如果颜色是按照十六进制的格式提供的
        # - 6位十六进制 ： #RRGGBB ，如 #FF0000 （红色）
        # - 3位十六进制 ： #RGB ，如 #F00 （等同于 #FF0000 ）
        # - 8位十六进制 ： #RRGGBBAA ，包含透明度，如 #FF0000FF，暂不支持
        if color_str.startswith("#"):
            hex_color = color_str[1:]  # 把#去掉
            # 处理三位颜色
            if len(hex_color) == 3:
                hex_color = "".join([c * 2 for c in hex_color])
            # 处理六位颜色
            if len(hex_color) == 6:
                try:
                    # int的第2个参数表明进制
                    r = int(hex_color[0:2], 16)
                    g = int(hex_color[2:4], 16)
                    b = int(hex_color[4:6], 16)
                    return RGBColor(r, g, b)
                except ValueError:
                    pass

        # 2. 如果颜色是按照rgb(r,g,b)的格式提供的
        # 使用正则表达式匹配RGB格式的颜色字符串
        # 匹配的格式为: rgb(r,g,b), 其中r,g,b为0-255的整数
        # \s* 表示匹配0个或多个空白字符
        # (\d+) 表示匹配1个或多个数字,并捕获到组中
        rgb_match = re.match(
            r"rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", color_str
        )

        # 如果匹配成功
        if rgb_match:
            # 使用map将匹配到的字符串转为整数
            # groups()返回所有捕获组的元组
            r, g, b = map(int, rgb_match.groups())
            # 返回RGBColor对象
            return RGBColor(r, g, b)

        # 3. 如果颜色是按照命名颜色的格式提供的
        color_map = {
            "black": RGBColor(0, 0, 0),
            "white": RGBColor(255, 255, 255),
            "red": RGBColor(255, 0, 0),
            "green": RGBColor(0, 128, 0),
            "blue": RGBColor(0, 0, 255),
            "gray": RGBColor(128, 128, 128),
            "grey": RGBColor(128, 128, 128),
        }

        return color_map.get(color_str)


class BackgroundHelper:
    """背景处理助手

    Attributes:
        config (SlideConfig): 幻灯片配置对象，包含幻灯片的相关配置信息。
    """

    def __init__(self, config: SlideConfig):
        self.config = config
        # 背景处理助手
        self.color_helper = ColorHelper()

    def create_background_image(self, background_style: str) -> Optional[str]:
        """创建背景图片

        Args:
            background_style (str): 背景样式字符串，支持渐变、纯色、图片等。
        Returns:
            Optional[str]: 背景图片的路径，如果无法创建则返回None。
        """
        if not background_style:
            return None

        # 处理渐变背景
        if "gradient" in background_style.lower():
            return self._create_gradient_background(background_style)

        # 处理纯色背景
        color = ColorHelper.parse_color(background_style)
        if color:
            return self._create_solid_background(color)

        return None

    def _parse_background(self, style_dict: dict) -> Optional[str]:
        """解析背景样式"""
        background = style_dict.get("background", "")
        background_color = style_dict.get("background-color", "")

        # 处理渐变背景
        if "linear-gradient" in background:
            return self._create_gradient_background(background)

        # 处理纯色背景
        if background_color and background_color != "transparent":
            color = self.color_helper.parse_color(background_color)
            if color:
                return self._create_solid_background(color)

        # 处理background简写属性中的颜色
        if (
            background
            and not background.startswith("url")
            and "gradient" not in background
        ):
            # 提取颜色值（可能包含其他属性）
            color_match = re.search(
                r"#[0-9a-fA-F]{3,6}|rgb\([^)]+\)|rgba\([^)]+\)|[a-zA-Z]+", background
            )
            if color_match:
                color_value = color_match.group()
                color = self.color_helper.parse_color(color_value)
                if color:
                    return self._create_solid_background(color)

        return None

    def _create_gradient_background(self, gradient_str: str) -> Optional[str]:
        """创建渐变背景"""
        try:
            # 解析linear-gradient，支持更多格式
            # 匹配颜色值：十六进制、rgb、rgba、颜色名称
            color_pattern = r"#[0-9a-fA-F]{3,6}|rgb\([^)]+\)|rgba\([^)]+\)|\b(?:red|blue|green|white|black|gray|grey|yellow|orange|purple|pink|brown|cyan|magenta)\b"
            colors = re.findall(color_pattern, gradient_str, re.IGNORECASE)

            if len(colors) >= 2:
                start_color = self.color_helper.parse_color(colors[0])
                end_color = self.color_helper.parse_color(colors[-1])
            else:
                # 如果没有找到足够的颜色，使用默认渐变
                start_color = self.color_helper.parse_color("#667eea")
                end_color = self.color_helper.parse_color("#764ba2")

            if not start_color or not end_color:
                return None

            # 创建渐变图片
            img = Image.new("RGB", (self.config.width_px, self.config.height_px))
            draw = ImageDraw.Draw(img)

            # 解析起始和结束颜色的RGB值
            start_str = str(start_color)
            end_str = str(end_color)

            if len(start_str) == 6 and len(end_str) == 6:
                start_r = int(start_str[0:2], 16)
                start_g = int(start_str[2:4], 16)
                start_b = int(start_str[4:6], 16)

                end_r = int(end_str[0:2], 16)
                end_g = int(end_str[2:4], 16)
                end_b = int(end_str[4:6], 16)

                # 检查渐变方向
                if "135deg" in gradient_str or "to bottom right" in gradient_str:
                    # 对角线渐变
                    for y in range(self.config.height_px):
                        for x in range(self.config.width_px):
                            ratio = (x + y) / (
                                self.config.width_px + self.config.height_px
                            )
                            r = int(start_r * (1 - ratio) + end_r * ratio)
                            g = int(start_g * (1 - ratio) + end_g * ratio)
                            b = int(start_b * (1 - ratio) + end_b * ratio)
                            draw.point((x, y), fill=(r, g, b))
                else:
                    # 垂直渐变（默认）
                    for y in range(self.config.height_px):
                        ratio = y / self.config.height_px
                        r = int(start_r * (1 - ratio) + end_r * ratio)
                        g = int(start_g * (1 - ratio) + end_g * ratio)
                        b = int(start_b * (1 - ratio) + end_b * ratio)
                        draw.line([(0, y), (self.config.width_px, y)], fill=(r, g, b))

            # 保存图片
            bg_filename = f"gradient_bg_{abs(hash(gradient_str)) % 10000}.png"
            img.save(bg_filename, "PNG")
            return bg_filename

        except Exception as e:
            print(f"创建渐变背景失败: {e}")
            return None

    def _create_solid_background(self, color: RGBColor) -> Optional[str]:
        """创建纯色背景"""
        try:
            # RGBColor对象实际上是一个字符串形式的十六进制颜色值
            color_str = str(color)
            if len(color_str) == 6:  # RRGGBB格式
                r = int(color_str[0:2], 16)
                g = int(color_str[2:4], 16)
                b = int(color_str[4:6], 16)
            else:
                # 默认白色背景
                r, g, b = 255, 255, 255

            img = Image.new(
                "RGB", (self.config.width_px, self.config.height_px), (r, g, b)
            )

            bg_filename = f"solid_bg_{r}_{g}_{b}.png"
            img.save(bg_filename, "PNG")
            return bg_filename

        except Exception as e:
            print(f"创建纯色背景失败: {e}")
            return None


class HTML2PPTXConverter:
    """HTML转PPT转换器"""

    def __init__(self, config: SlideConfig = None):
        self.config = config or SlideConfig()
        self.background_helper = BackgroundHelper(self.config)
        self.temp_files = []

    def convert(
        self, html_file: str, css_file: str = None, output_file: str = None
    ) -> str:
        """转换HTML文件为PPT"""
        print(f"开始转换: {html_file}")

        # 读取HTML文件
        with open(html_file, "r", encoding="utf-8") as f:
            html_content = f.read()

        # 读取CSS文件
        css_styles = {}
        if css_file and os.path.exists(css_file):
            print(f"加载CSS文件: {css_file}")
            css_styles = self._parse_css_file(css_file)

        # 解析HTML
        soup = BeautifulSoup(html_content, "html.parser")

        # 创建PPT
        prs = Presentation()
        prs.slide_width = Inches(self.config.width_inches)
        prs.slide_height = Inches(self.config.height_inches)

        # 查找幻灯片容器
        slide_containers: List[Tag] = soup.find_all("div", class_="slide-container")
        print(f"找到 {len(slide_containers)} 个幻灯片容器")

        if not slide_containers:
            # 如果没有找到slide-container，将整个body作为一个幻灯片
            body = soup.find("body")
            if body:
                slide_containers = [body]
                print("使用body作为幻灯片容器")

        # 处理每个幻灯片，每个container都是一张幻灯片
        # 一页一页生成
        for i, container in enumerate(slide_containers):
            print(f"处理第 {i+1} 个幻灯片")
            self._create_slide(prs, container, css_styles)

        # 保存PPT
        if not output_file:  # 如果没有指定输出文件名
            import time  # 导入time模块，根据当前时间生成文件名

            base_name = os.path.splitext(html_file)[0]
            timestamp = int(time.time())
            output_file = f"{base_name}_fixed_{timestamp}.pptx"

        prs.save(output_file)
        print(f"转换完成: {output_file}")

        # 清理临时文件
        self._cleanup_temp_files()

        # 清理临时文件
        self._cleanup_temp_files()

        return output_file

    def _parse_css_file(self, css_file: str) -> Dict[str, Dict[str, str]]:
        """解析CSS文件

        Args:
            css_file (str): CSS文件路径

        Returns:
            Dict[str, Dict[str, str]]: 解析后的CSS样式字典，例如：
            {
                ".title-slide": {
                    "background": "linear-gradient(135deg, #f8f9fa 0%, #ffffff 50%, #f0f8ff 100%)",
                    "color": "#333333",
                    "text-align": "center",
                    "justify-content": "center",
                    "align-items": "center",
                    "border-top": "4px solid #2c5aa0",
                    "border-bottom": "2px solid #2c5aa0",
                    "position": "relative",
                },
                ".title-slide::after": {
                    "content": "''",
                    "position": "absolute",
                    "bottom": "60px",
                    "left": "50%",
                    "transform": "translateX(-50%)",
                    "width": "200px",
                    "height": "2px",
                    "background": "linear-gradient(90deg, transparent 0%, #2c5aa0 50%, transparent 100%)",
                },
                ……
            }
        """
        styles = {}

        try:
            with open(css_file, "r", encoding="utf-8") as f:
                css_content = f.read()

            # 简单的CSS解析
            # 移除注释
            css_content = re.sub(r"/\*.*?\*/", "", css_content, flags=re.DOTALL)

            # 解析规则:
            # 1. ([^{}]+) - 匹配选择器部分
            #   - [^{}]+ 表示匹配除了{}以外的任意字符一次或多次
            #   - () 表示捕获这个分组
            # 2. \{ - 匹配左花括号
            # 3. ([^{}]+) - 匹配CSS属性部分
            #   - 同样使用[^{}]+匹配除了{}以外的所有字符
            # 4. \} - 匹配右花括号
            rules = re.findall(r"([^{}]+)\{([^{}]+)\}", css_content)

            for selector, properties in rules:
                selector = selector.strip()
                props = {}

                for prop in properties.split(";"):
                    if ":" in prop:
                        key, value = prop.split(":", 1)  # 第二个参数为1，确保只分割一次
                        props[key.strip()] = value.strip()

                if props:
                    styles[selector] = props

        except Exception as e:
            print(f"CSS解析错误: {e}")

        return styles

    def _create_slide(
        self, prs: Presentation, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> None:
        """创建单个幻灯片

        Args:
            prs (Presentation): PPTX对象，用于添加幻灯片
            container (Tag): 幻灯片容器元素，从HTML解析得到
            css_styles (Dict[str, Dict[str, str]]): 解析后的CSS样式字典
        """
        # 使用空白布局
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 设置背景
        self._set_slide_background(slide, container, css_styles)

        # 提取并添加内容
        self._add_slide_content(slide, container, css_styles)

    def _set_slide_background(
        self, slide, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> None:
        """智能设置幻灯片背景"""
        # 获取背景样式和装饰元素
        background_info = self._analyze_background_style(container, css_styles)

        if background_info["needs_decoration"]:
            # 需要装饰元素，生成背景图片
            bg_image_path = self._create_decorated_background(
                background_info, css_styles
            )
            # 检查背景图片是否存在，如存在，才添加
            if bg_image_path and os.path.exists(bg_image_path):
                try:
                    slide.shapes.add_picture(
                        bg_image_path,
                        0,
                        0,
                        Inches(self.config.width_inches),
                        Inches(self.config.height_inches),
                    )
                    self.temp_files.append(bg_image_path)
                    print(f"添加装饰背景图片: {bg_image_path}")
                except Exception as e:
                    print(f"设置装饰背景失败: {e}")
        else:
            # 纯色背景，直接设置PPT背景色
            self._set_solid_background_color(slide, background_info["background_color"])
            print(f"设置纯色背景: {background_info['background_color']}")

    def _analyze_background_style(
        self, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> Dict[str, Any]:
        """分析背景样式和装饰需求

        Args:
            container (Tag): 幻灯片容器元素，从HTML解析得到
            css_styles (Dict[str, Dict[str, str]]): 解析后的CSS样式字典

        Returns:
            Dict[str, Any]: 包含背景颜色、是否需要装饰、装饰信息等的字典，目前包括四个键-值对：
                - background_color (str): 背景颜色
                - needs_decoration (bool): 是否需要装饰
                - decoration_info (Dict[str, str]): 装饰信息
                - container_classes (List[str]): 容器的类名列表
        """
        # 检查容器的class属性
        container_classes = container.get("class", [])
        if isinstance(container_classes, str):
            container_classes = [container_classes]

        background_color = "#ffffff"  # 默认白色
        needs_decoration = False  # 是否需要装饰
        decoration_info = {}

        # 从特定类样式获取背景和装饰信息
        for class_name in container_classes:
            # 从HTML的class属性中获取类名，转为css中的类选择器，例如：.title-slide，要加“.”
            class_selector = f".{class_name}"
            if class_selector in css_styles:
                class_styles = css_styles[class_selector]

                # 获取背景色
                bg_style = class_styles.get("background") or class_styles.get(
                    "background-color"
                )
                if bg_style:
                    background_color = bg_style

                # 检查是否有装饰元素，HTML的上边框或者下边框
                if class_styles.get("border-top"):
                    needs_decoration = True
                    decoration_info["border_top"] = class_styles["border-top"]

                if class_styles.get("border-bottom"):
                    needs_decoration = True
                    decoration_info["border_bottom"] = class_styles["border-bottom"]

                # 检查是否有渐变背景
                if bg_style and (
                    "gradient" in bg_style.lower()
                    or "linear-gradient" in bg_style.lower()
                ):
                    needs_decoration = True
                    decoration_info["gradient"] = bg_style

                # 检查伪元素装饰（通过类名推断），一般标题页会有
                # if class_name == "title-slide":
                #     needs_decoration = True
                #     decoration_info["title_decoration"] = True

        # 从slide-container样式获取背景
        # if ".slide-container" in css_styles:
        #     container_styles = css_styles[".slide-container"]
        #     if not background_color or background_color == "#ffffff":
        #         bg_style = container_styles.get("background") or container_styles.get(
        #             "background-color"
        #         )
        #         if bg_style:
        #             background_color = bg_style

        return {
            "background_color": background_color,
            "needs_decoration": needs_decoration,
            "decoration_info": decoration_info,
            "container_classes": container_classes,
        }

    def _create_decorated_background(
        self, background_info: Dict, css_styles: Dict[str, Dict[str, str]]
    ) -> Optional[str]:
        """创建带装饰的背景图片"""
        try:
            from PIL import Image, ImageDraw

            # 创建基础背景
            img = Image.new("RGB", (self.config.width_px, self.config.height_px))
            draw = ImageDraw.Draw(img)

            # 设置基础背景色
            bg_color = self._parse_color_to_rgb(background_info["background_color"])
            img.paste(bg_color, (0, 0, self.config.width_px, self.config.height_px))

            # 添加装饰元素
            decoration_info = background_info["decoration_info"]

            # 添加顶部边框
            if "border_top" in decoration_info:
                border_info = self._parse_border(decoration_info["border_top"])
                if border_info:
                    border_color = self._parse_color_to_rgb(border_info["color"])
                    border_width = border_info["width"]
                    draw.rectangle(
                        [(0, 0), (self.config.width_px, border_width)],
                        fill=border_color,
                    )

            # 添加底部边框
            if "border_bottom" in decoration_info:
                border_info = self._parse_border(decoration_info["border_bottom"])
                if border_info:
                    border_color = self._parse_color_to_rgb(border_info["color"])
                    border_width = border_info["width"]
                    draw.rectangle(
                        [
                            (0, self.config.height_px - border_width),
                            (self.config.width_px, self.config.height_px),
                        ],
                        fill=border_color,
                    )

            # 处理渐变背景
            if "gradient" in decoration_info:
                # 使用现有的渐变背景生成器
                gradient_bg = self.background_helper.create_background_image(
                    decoration_info["gradient"]
                )
                if gradient_bg:
                    # 在渐变背景上添加其他装饰元素
                    return self._add_decorations_to_gradient(
                        gradient_bg, decoration_info
                    )

            # 添加标题页特殊装饰
            if decoration_info.get("title_decoration"):
                self._add_title_decorations(draw, decoration_info)

            # 保存装饰背景图片
            bg_filename = f"decorated_bg_{abs(hash(str(decoration_info))) % 10000}.png"
            img.save(bg_filename, "PNG")
            return bg_filename

        except Exception as e:
            print(f"创建装饰背景失败: {e}")
            # 回退到原有方法
            return self.background_helper.create_background_image(
                background_info["background_color"]
            )

    def _parse_color_to_rgb(self, color_str: str) -> tuple:
        """解析颜色字符串为RGB元组

        Args:
            color_str (str): 颜色字符串，支持十六进制、颜色名称等。
        Returns:
            tuple: 包含RGB值的元组，例如：(255, 0, 0)。
        """
        color_str = color_str.strip().lower()

        if color_str.startswith("#"):
            color_str = color_str[1:]
            if len(color_str) == 6:
                return (
                    int(color_str[0:2], 16),
                    int(color_str[2:4], 16),
                    int(color_str[4:6], 16),
                )
            elif len(color_str) == 3:
                return (
                    int(color_str[0] * 2, 16),
                    int(color_str[1] * 2, 16),
                    int(color_str[2] * 2, 16),
                )

        # 常见颜色名称
        color_map = {
            "white": (255, 255, 255),
            "black": (0, 0, 0),
            "red": (255, 0, 0),
            "green": (0, 128, 0),
            "blue": (0, 0, 255),
            "yellow": (255, 255, 0),
            "cyan": (0, 255, 255),
            "magenta": (255, 0, 255),
        }

        return color_map.get(color_str, (255, 255, 255))  # 默认白色

    def _parse_border(self, border_str: str) -> Optional[Dict]:
        """解析边框样式字符串"""
        try:
            # 解析类似 "4px solid #2c5aa0" 的边框样式
            parts = border_str.strip().split()
            if len(parts) >= 3:
                width_str = parts[0]
                style = parts[1]
                color = parts[2]

                # 解析宽度（转换px到实际像素）
                width = 4  # 默认宽度
                if width_str.endswith("px"):
                    width = int(width_str[:-2])

                return {"width": width, "style": style, "color": color}
        except Exception as e:
            print(f"解析边框样式失败: {e}")

        return None

    def _set_solid_background_color(self, slide, color_str: str) -> None:
        """设置纯色背景"""
        try:
            # 解析颜色
            color = ColorHelper.parse_color(color_str)
            if color:
                # 设置幻灯片背景色
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = color
                print(f"设置幻灯片背景色: {color_str}")
            else:
                print(f"无法解析颜色: {color_str}，使用默认白色背景")
        except Exception as e:
            print(f"设置背景色失败: {e}")

    def _add_decorations_to_gradient(
        self, gradient_bg_path: str, decoration_info: Dict
    ) -> str:
        """在渐变背景上添加装饰元素

        Args:
            gradient_bg_path (str): 渐变背景图片路径。
            decoration_info (Dict): 包含装饰信息的字典。
        Returns:
            str: 增强后的背景图片路径。
        """
        try:
            from PIL import Image, ImageDraw

            # 打开渐变背景图片
            img = Image.open(gradient_bg_path)
            draw = ImageDraw.Draw(img)

            # 添加边框装饰
            if "border_top" in decoration_info:
                border_info = self._parse_border(decoration_info["border_top"])
                if border_info:
                    border_color = self._parse_color_to_rgb(border_info["color"])
                    border_width = border_info["width"]
                    draw.rectangle(
                        [(0, 0), (self.config.width_px, border_width)],
                        fill=border_color,
                    )

            if "border_bottom" in decoration_info:
                border_info = self._parse_border(decoration_info["border_bottom"])
                if border_info:
                    border_color = self._parse_color_to_rgb(border_info["color"])
                    border_width = border_info["width"]
                    draw.rectangle(
                        [
                            (0, self.config.height_px - border_width),
                            (self.config.width_px, self.config.height_px),
                        ],
                        fill=border_color,
                    )

            # 添加标题页装饰线
            if decoration_info.get("title_decoration"):
                self._add_title_decorations(draw, decoration_info)

            # 保存增强的背景图片
            enhanced_bg_filename = (
                f"enhanced_bg_{abs(hash(str(decoration_info))) % 10000}.png"
            )
            img.save(enhanced_bg_filename, "PNG")

            # 删除原始渐变背景文件
            if os.path.exists(gradient_bg_path):
                os.remove(gradient_bg_path)

            return enhanced_bg_filename

        except Exception as e:
            print(f"添加装饰元素失败: {e}")
            return gradient_bg_path

    def _add_title_decorations(self, draw, decoration_info: Dict) -> None:
        """添加标题页特殊装饰元素"""
        try:
            # 添加底部装饰线（模拟CSS的::after伪元素）
            line_width = 200
            line_height = 2
            line_y = self.config.height_px - 120  # 距离底部60px转换为像素
            line_x = (self.config.width_px - line_width) // 2

            # 创建渐变装饰线
            line_color = self._parse_color_to_rgb("#2c5aa0")

            # 绘制渐变装饰线
            for i in range(line_width):
                # 计算透明度渐变
                distance_from_center = abs(i - line_width // 2)
                alpha_factor = 1.0 - (distance_from_center / (line_width // 2))
                alpha_factor = max(0.0, min(1.0, alpha_factor))

                # 应用透明度（通过调整颜色亮度模拟）
                adjusted_color = tuple(
                    int(c * alpha_factor + 255 * (1 - alpha_factor)) for c in line_color
                )

                draw.rectangle(
                    [(line_x + i, line_y), (line_x + i + 1, line_y + line_height)],
                    fill=adjusted_color,
                )

        except Exception as e:
            print(f"添加标题装饰失败: {e}")

    def _add_slide_content(
        self, slide, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> None:
        """添加幻灯片内容"""
        # 检查是否为标题页，调整初始位置
        is_title_slide = "title-slide" in container.get("class", [])
        if is_title_slide:
            current_y = self.config.padding_top
        else:
            current_y = (
                self.config.padding_top + 0.1
            )  # 进一步减少初始间距，让内容更靠上

        # 检查是否有复杂效果需要转换为图片
        if self._has_complex_effects(container, css_styles):
            self._add_complex_content_as_image(slide, container, css_styles)
            return

        # 只获取容器的直接子元素，避免重复处理
        content_elements = []
        for child in container.children:
            if hasattr(child, "name") and child.name in [
                "h1",
                "h2",
                "h3",
                "p",
                "ul",
                "ol",
                "div",
            ]:
                content_elements.append(child)

        for element in content_elements:
            text = element.get_text().strip()
            if not text:
                continue

            # 根据元素类型和上下文确定文本样式
            if element.name == "h1":
                if is_title_slide:
                    current_y = self._add_styled_text(
                        slide, text, current_y, "main_title", css_styles, element
                    )
                else:
                    current_y = self._add_styled_text(
                        slide, text, current_y, "slide_title", css_styles, element
                    )
            elif element.name == "h2":
                if is_title_slide:
                    current_y = self._add_styled_text(
                        slide, text, current_y, "subtitle", css_styles, element
                    )
                else:
                    current_y = self._add_styled_text(
                        slide, text, current_y, "heading", css_styles, element
                    )
            elif element.name == "h3":
                current_y = self._add_styled_text(
                    slide, text, current_y, "heading", css_styles, element
                )
            elif element.name == "p":
                # 检查是否为作者信息或特殊类型
                if is_title_slide and (
                    "author" in element.get("class", [])
                    or any(
                        keyword in text.lower()
                        for keyword in ["作者", "演讲者", "报告人"]
                    )
                ):
                    current_y = self._add_styled_text(
                        slide, text, current_y, "author", css_styles, element
                    )
                else:
                    current_y = self._add_styled_text(
                        slide, text, current_y, "body_text", css_styles, element
                    )
            elif element.name in ["ul", "ol"]:
                current_y = self._add_list(slide, element, current_y, css_styles)
            elif element.name == "div":
                # 处理div中的内容
                current_y = self._add_div_content(slide, element, current_y, css_styles)

            # 检查是否超出幻灯片边界
            if current_y > self.config.height_inches - self.config.padding_bottom:
                break

    def _add_styled_text(
        self,
        slide,
        text: str,
        y_position: float,
        style_name: str,
        css_styles: Dict[str, Dict[str, str]],
        element: Tag = None,
    ) -> float:
        """使用指定样式添加文本，支持内联样式"""
        style = self.config.text_styles.get(
            style_name, self.config.text_styles["body_text"]
        )

        width = (
            self.config.width_inches
            - self.config.padding_left
            - self.config.padding_right
        )

        # 根据CSS规格精确计算文本框高度
        chars_per_line = max(40, int(width * 30))  # 更保守的字符数估算
        estimated_lines = max(1, len(text) // chars_per_line + 1)  # 增加一行缓冲
        line_height = style.font_size / 72 * 1.3  # 增加行高倍数
        height = max(0.6, estimated_lines * line_height + 0.3)  # 增加额外高度，避免截断

        textbox = slide.shapes.add_textbox(
            Inches(self.config.padding_left),
            Inches(y_position),
            Inches(width),
            Inches(height),
        )

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.05)  # 按CSS规格减少边距
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.05)  # 减少上下边距
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        p = text_frame.paragraphs[0]

        # 如果有HTML元素，解析内联样式
        if element is not None:
            self._add_text_with_inline_styles(p, element, style, css_styles)
        else:
            p.text = text
            # 设置默认字体样式
            font = p.runs[0].font
            font.name = self.config.default_font
            font.size = Pt(style.font_size)
            font.bold = style.bold
            font.italic = style.italic

            # 设置字体颜色
            color_str = style.font_color
            element_styles = css_styles.get(
                self._get_css_selector_for_style(style_name), {}
            )
            if "color" in element_styles:
                color_str = element_styles["color"]

            color = ColorHelper.parse_color(color_str)
            if color:
                font.color.rgb = color

        # 设置对齐方式
        alignment_map = {
            "left": PP_PARAGRAPH_ALIGNMENT.LEFT,
            "center": PP_PARAGRAPH_ALIGNMENT.CENTER,
            "right": PP_PARAGRAPH_ALIGNMENT.RIGHT,
            "justify": PP_PARAGRAPH_ALIGNMENT.JUSTIFY,
        }
        p.alignment = alignment_map.get(style.alignment, PP_PARAGRAPH_ALIGNMENT.LEFT)

        # 按CSS规格设置段落间距
        p.line_spacing = style.line_spacing  # 使用CSS标准行间距
        p.space_before = Pt(style.space_before)  # 使用CSS标准段前间距
        p.space_after = Pt(style.space_after)  # 使用CSS标准段后间距

        # 为slide_title添加下划线效果
        if style_name == "slide_title":
            self._add_title_underline(slide, y_position, width)

        return y_position + height + 0.25  # 增加元素间距，避免文本框重叠

    def _add_text_with_inline_styles(
        self,
        paragraph,
        element: Tag,
        base_style: "TextStyle",
        css_styles: Dict[str, Dict[str, str]],
    ) -> None:
        """为段落添加带有内联样式的文本"""
        # 不清空段落，因为可能已经有前缀内容（如列表标记）

        def process_element(elem):
            """递归处理元素及其子元素"""
            if elem.name is None:  # 文本节点
                text_content = str(elem)
                if text_content:  # 不去除空格，保持原始格式
                    run = paragraph.add_run()
                    run.text = text_content
                    self._apply_base_style_to_run(run, base_style)
                    return run
            else:
                # HTML元素
                if elem.name == "strong":
                    # 处理strong标签 - 应用蓝色加粗样式
                    text_content = elem.get_text()
                    if text_content:
                        run = paragraph.add_run()
                        run.text = text_content
                        self._apply_base_style_to_run(run, base_style)

                        # 应用strong样式：蓝色 + 加粗
                        run.font.bold = True
                        run.font.color.rgb = ColorHelper.parse_color(
                            "#2c5aa0"
                        )  # CSS中定义的蓝色
                        return run
                elif elem.name in ["em", "i"]:
                    # 处理斜体标签
                    text_content = elem.get_text()
                    if text_content:
                        run = paragraph.add_run()
                        run.text = text_content
                        self._apply_base_style_to_run(run, base_style)
                        run.font.italic = True
                        return run
                elif elem.name in ["b"]:
                    # 处理粗体标签
                    text_content = elem.get_text()
                    if text_content:
                        run = paragraph.add_run()
                        run.text = text_content
                        self._apply_base_style_to_run(run, base_style)
                        run.font.bold = True
                        return run
                else:
                    # 其他元素，递归处理子元素
                    for child in elem.children:
                        process_element(child)
            return None

        # 处理元素的所有子内容
        for child in element.children:
            process_element(child)

        # 如果段落中没有任何run（除了可能的前缀），添加默认文本
        if len(paragraph.runs) <= 1:  # 考虑可能已有的前缀run
            text_content = element.get_text().strip()
            if text_content:
                run = paragraph.add_run()
                run.text = text_content
                self._apply_base_style_to_run(run, base_style)

    def _apply_base_style_to_run(self, run, base_style: "TextStyle") -> None:
        """为run应用基础样式"""
        font = run.font
        font.name = self.config.default_font
        font.size = Pt(base_style.font_size)
        if not font.bold:  # 只有在没有被内联样式设置为粗体时才应用基础样式
            font.bold = base_style.bold
        if not font.italic:  # 只有在没有被内联样式设置为斜体时才应用基础样式
            font.italic = base_style.italic

        # 设置默认颜色（如果没有被内联样式覆盖）
        if not hasattr(font.color, "_color_val") or font.color.rgb is None:
            color = ColorHelper.parse_color(base_style.font_color)
            if color:
                font.color.rgb = color

    def _add_title_underline(self, slide, y_position: float, width: float) -> None:
        """为标题添加下划线效果"""
        try:
            from pptx.shapes.connector import Connector
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.util import Inches
            from pptx.dml.color import RGBColor

            # 计算下划线位置（模拟CSS border-bottom效果）
            underline_y = y_position + 0.65  # 进一步增加标题与下划线的距离
            underline_width = width * 0.95  # 下划线宽度为文本宽度的95%

            # 添加下划线形状
            line = slide.shapes.add_connector(
                connector_type=1,  # 直线
                begin_x=Inches(self.config.padding_left + width * 0.025),  # 居中对齐
                begin_y=Inches(underline_y),
                end_x=Inches(self.config.padding_left + width * 0.975),
                end_y=Inches(underline_y),
            )

            # 设置线条样式
            line.line.color.rgb = RGBColor(44, 90, 160)  # #2c5aa0
            line.line.width = Inches(0.02)  # 3px转换为英寸

        except Exception as e:
            print(f"添加标题下划线失败: {e}")

    def _get_css_selector_for_style(self, style_name: str) -> str:
        """根据样式名称获取对应的CSS选择器"""
        selector_map = {
            "main_title": "h1",
            "slide_title": "h1",
            "subtitle": "h2",
            "heading": "h2",
            "body_text": "p",
            "author": "p",
            "caption": "p",
        }
        return selector_map.get(style_name, "p")

    def _has_complex_effects(
        self, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> bool:
        """检测是否有复杂的视觉效果"""
        # 检查容器类名
        container_classes = container.get("class", [])
        if isinstance(container_classes, str):
            container_classes = [container_classes]

        # 复杂效果的类名列表 - 更精确的匹配
        complex_classes = [
            "chart-container",
            "chart-item",
            "chart-bar",
            "code-container",
            "code-block",
            "insights-container",
            "insight-card",
            "stat-box",
            "two-column-slide",
            "feature-list",
            "data-container",
        ]

        # 检查是否包含复杂效果类名
        for class_name in container_classes:
            if class_name in complex_classes:
                return True

        # 检查子元素是否有复杂效果
        for complex_class in complex_classes:
            if container.find(class_=complex_class):
                return True

        # 检查是否有多列布局或复杂结构
        if container.find_all(class_=["column", "col", "grid-item"]):
            return True

        # 检查是否有图表或数据可视化元素
        if container.find_all(["canvas", "svg"]) or container.find(
            class_=lambda x: x and ("chart" in x or "graph" in x)
        ):
            return True

        return False

    def _add_complex_content_as_image(
        self, slide, container: Tag, css_styles: Dict[str, Dict[str, str]]
    ) -> None:
        """将复杂内容转换为图片添加到幻灯片"""
        try:
            # 这里可以使用selenium或其他工具来截图
            # 暂时用文本框显示提示信息
            width = (
                self.config.width_inches
                - self.config.padding_left
                - self.config.padding_right
            )
            height = 2.0

            textbox = slide.shapes.add_textbox(
                Inches(self.config.padding_left),
                Inches(self.config.padding_top + 2),
                Inches(width),
                Inches(height),
            )

            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.margin_left = Inches(0.3)
            text_frame.margin_right = Inches(0.3)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            p.text = f"复杂效果内容: {container.get_text()[:100]}..."
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

            # 设置字体样式
            font = p.runs[0].font
            font.name = self.config.default_font
            font.size = Pt(18)
            font.italic = True

            print(f"检测到复杂效果，已转换为文本提示")

        except Exception as e:
            print(f"处理复杂内容失败: {e}")

    def _add_list(
        self,
        slide,
        list_elem: Tag,
        y_position: float,
        css_styles: Dict[str, Dict[str, str]],
    ) -> float:
        """添加列表"""
        style = self.config.text_styles["list_item"]
        width = (
            self.config.width_inches
            - self.config.padding_left
            - self.config.padding_right
            - 0.3
        )

        # 获取列表项
        list_items = list_elem.find_all("li")
        if not list_items:
            return y_position

        # 更精确的高度估算
        item_count = len([item for item in list_items if item.get_text().strip()])
        avg_text_length = sum(
            len(item.get_text().strip())
            for item in list_items
            if item.get_text().strip()
        ) / max(1, item_count)
        chars_per_line = max(30, int(width * 25))
        lines_per_item = max(1, avg_text_length // chars_per_line)
        line_height = style.font_size / 72 * 1.2
        height = max(1.0, item_count * lines_per_item * line_height + 0.4)

        textbox = slide.shapes.add_textbox(
            Inches(self.config.padding_left + 0.3),
            Inches(y_position),
            Inches(width),
            Inches(height),
        )

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # 添加列表项
        is_ordered = list_elem.name == "ol"

        for i, li in enumerate(list_items):
            text = li.get_text().strip()
            if not text:
                continue

            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # 添加列表标记
            if is_ordered:
                prefix = f"{i + 1}. "
            else:
                prefix = "• "

            # 清空段落并添加前缀
            p.clear()
            prefix_run = p.add_run()
            prefix_run.text = prefix
            self._apply_base_style_to_run(prefix_run, style)

            # 处理列表项的内联样式
            self._add_text_with_inline_styles(p, li, style, css_styles)

            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            p.line_spacing = style.line_spacing
            p.space_before = Pt(style.space_before)
            p.space_after = Pt(style.space_after)

        return y_position + height + 0.25  # 增加列表后的间距，避免重叠

    def _add_div_content(
        self,
        slide,
        div_element: Tag,
        y_position: float,
        css_styles: Dict[str, Dict[str, str]],
    ) -> float:
        """处理div中的内容"""
        current_y = y_position

        # 获取div中的所有子元素（包括深层嵌套的）
        content_elements = div_element.find_all(
            ["h1", "h2", "h3", "p", "ul", "ol"], recursive=True
        )

        # 如果没有找到结构化内容，将整个div作为段落处理
        if not content_elements:
            text = div_element.get_text().strip()
            if text:
                current_y = self._add_styled_text(
                    slide, text, current_y, "body_text", css_styles, div_element
                )
            return current_y

        # 处理找到的结构化内容
        for element in content_elements:
            text = element.get_text().strip()
            if not text:
                continue

            if element.name == "h1":
                current_y = self._add_styled_text(
                    slide, text, current_y, "slide_title", css_styles, element
                )
            elif element.name == "h2":
                current_y = self._add_styled_text(
                    slide, text, current_y, "heading", css_styles, element
                )
            elif element.name == "h3":
                current_y = self._add_styled_text(
                    slide, text, current_y, "heading", css_styles, element
                )
            elif element.name == "p":
                current_y = self._add_styled_text(
                    slide, text, current_y, "body_text", css_styles, element
                )
            elif element.name in ["ul", "ol"]:
                current_y = self._add_list(slide, element, current_y, css_styles)

        return current_y

    def _cleanup_temp_files(self) -> None:
        """清理临时文件"""
        for file_path in self.temp_files:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    print(f"清理临时文件: {file_path}")
            except Exception as e:
                print(f"清理文件失败 {file_path}: {e}")

        self.temp_files.clear()


def main():
    """主函数"""
    converter = FixedHTML2PPTXConverter()

    # 转换ppt-demo
    ppt_demo_html = "ppt-demo/ppt-demo.html"
    ppt_demo_css = "ppt-demo/ppt-template.css"

    if os.path.exists(ppt_demo_html):
        output_file = converter.convert(ppt_demo_html, ppt_demo_css)
        print(f"\n转换完成: {output_file}")
    else:
        print(f"文件不存在: {ppt_demo_html}")

    # 转换complex-demo
    complex_demo_html = "complex-demo/complex-demo.html"
    complex_demo_css = "complex-demo/complex-template.css"

    if os.path.exists(complex_demo_html):
        output_file = converter.convert(complex_demo_html, complex_demo_css)
        print(f"\n转换完成: {output_file}")
    else:
        print(f"文件不存在: {complex_demo_html}")


if __name__ == "__main__":
    main()
