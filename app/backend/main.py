import asyncio
import json
import os
from typing import Dict, Any, Optional

from fastapi import FastAPI, WebSocket, WebSocketDisconnect, UploadFile, File, Form
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from starlette.websockets import WebSocketState

from ppt_generate.agents.ppt_agent import PPTAgent
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import PyPDF2

app = FastAPI(title="PPT Agent Backend", version="0.1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/health")
async def health() -> Dict[str, Any]:
    return {"status": "ok"}


def _save_upload(tmp_dir: str, uf: UploadFile) -> str:
    os.makedirs(tmp_dir, exist_ok=True)
    file_path = os.path.join(tmp_dir, uf.filename)
    with open(file_path, "wb") as f:
        f.write(uf.file.read())
    return file_path


@app.post("/upload")
async def upload(file: UploadFile = File(...)):
    try:
        saved_path = _save_upload("/tmp/ppt-agent-uploads", file)
        return {"success": True, "path": saved_path, "filename": file.filename}
    except Exception as e:
        return JSONResponse(status_code=500,
                            content={
                                "success": False,
                                "error": str(e)
                            })


def _extract_text_from_file(path: str) -> Dict[str, Any]:
    try:
        ext = os.path.splitext(path)[1].lower()
        if ext in [".txt", ".md", ".log"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return {"success": True, "text": f.read(), "type": ext}

        if ext in [".html", ".htm"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                return {
                    "success": True,
                    "text": soup.get_text("\n"),
                    "type": ext
                }

        if ext == ".pdf":
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                texts = []
                for i in range(len(reader.pages)):
                    try:
                        t = reader.pages[i].extract_text() or ""
                        if t.strip():
                            texts.append(f"=== 第 {i + 1} 页 ===\n{t}")
                    except Exception:
                        continue
                return {
                    "success": True,
                    "text": "\n\n".join(texts),
                    "type": ext
                }

        if ext in [".docx"]:
            doc = Document(path)
            paragraphs = [
                p.text for p in doc.paragraphs if p.text and p.text.strip()
            ]
            return {
                "success": True,
                "text": "\n".join(paragraphs),
                "type": ext
            }

        if ext in [".pptx"]:
            prs = Presentation(path)
            texts = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                if slide_text:
                    texts.append(f"=== 第 {slide_idx} 页 ===\n" +
                                 "\n".join(slide_text))
            return {"success": True, "text": "\n\n".join(texts), "type": ext}

        # 未知类型，按文本尝试
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return {
                "success": True,
                "text": f.read(),
                "type": ext or "unknown"
            }
    except Exception as e:
        return {"success": False, "error": str(e), "text": ""}


@app.post("/extract")
async def extract(path: str = Form(...)):
    if not os.path.exists(path):
        return JSONResponse(status_code=404,
                            content={
                                "success": False,
                                "error": "file not found"
                            })
    result = _extract_text_from_file(path)
    if not result.get("success"):
        return JSONResponse(status_code=500, content=result)
    return result


@app.websocket("/ws/generate")
async def ws_generate(websocket: WebSocket):
    await websocket.accept()

    agent: Optional[PPTAgent] = PPTAgent()

    async def _send_json_safe(payload: Dict[str, Any]):
        try:
            if websocket.client_state != WebSocketState.CONNECTED:
                return
            await websocket.send_text(json.dumps(payload))
        except Exception:
            pass

    def on_event(evt: Dict[str, Any]):
        asyncio.create_task(_send_json_safe(evt))

    try:
        while True:
            msg = await websocket.receive_text()
            data = json.loads(msg)
            action = data.get("action")

            # 兼容老协议：未带 action 的一次性执行
            if not action:
                query: str = data.get("query", "")
                reference_content: Optional[str] = data.get(
                    "reference_content")
                reference_path = data.get("reference_path")
                rethink: bool = bool(data.get("rethink", True))
                max_rethink_times: int = int(data.get("max_rethink_times", 3))

                if reference_path and not reference_content:
                    ext_result = _extract_text_from_file(reference_path)
                    if ext_result.get("success"):
                        reference_content = ext_result.get("text")
                        await _send_json_safe({
                            "stage":
                            "reference_loaded",
                            "type":
                            ext_result.get("type"),
                            "chars":
                            len(reference_content or ""),
                        })

                await agent.generate_ppt_outline(
                    query=query,
                    reference_content=reference_content,
                    on_event=on_event)
                await agent.generate_page_content(
                    outline=agent.ppt_info["outline"],
                    rethink=rethink,
                    max_rethink_times=max_rethink_times,
                    on_event=on_event,
                )
                await _send_json_safe({
                    "stage": "done",
                    "outline": agent.ppt_info.get("outline"),
                    "pages": agent.ppt_info.get("pages", []),
                })
                continue

            if action == "start_outline":
                query: str = data.get("query", "")
                reference_content: Optional[str] = data.get(
                    "reference_content")
                reference_path = data.get("reference_path")

                if reference_path and not reference_content:
                    ext_result = _extract_text_from_file(reference_path)
                    if ext_result.get("success"):
                        reference_content = ext_result.get("text")
                        await _send_json_safe({
                            "stage":
                            "reference_loaded",
                            "type":
                            ext_result.get("type"),
                            "chars":
                            len(reference_content or ""),
                        })

                await agent.generate_ppt_outline(
                    query=query,
                    reference_content=reference_content,
                    on_event=on_event)
                # 结束时 ppt_agent 会推送 outline_done
                continue

            if action == "start_content":
                outline: str = data.get("outline", "")
                rethink: bool = bool(data.get("rethink", True))
                max_rethink_times: int = int(data.get("max_rethink_times", 3))
                if not outline:
                    await _send_json_safe({
                        "stage": "error",
                        "error": "outline is required"
                    })
                    continue

                agent.ppt_info["outline"] = outline
                await agent.generate_page_content(
                    outline=outline,
                    rethink=rethink,
                    max_rethink_times=max_rethink_times,
                    on_event=on_event)
                await _send_json_safe({
                    "stage": "done",
                    "outline": agent.ppt_info.get("outline"),
                    "pages": agent.ppt_info.get("pages", []),
                })
                continue

            await _send_json_safe({
                "stage": "error",
                "error": f"unknown action: {action}"
            })

    except WebSocketDisconnect:
        pass
    except Exception as e:
        try:
            await _send_json_safe({"stage": "error", "error": str(e)})
        except Exception:
            pass
    finally:
        try:
            await websocket.close()
        except Exception:
            pass


if __name__ == "__main__":
    import uvicorn

    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
